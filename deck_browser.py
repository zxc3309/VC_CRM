import os
from io import BytesIO  # 移到最上面
import logging
from dotenv import load_dotenv
from utils.path_helper import PathHelper
import re
import asyncio
from bs4 import BeautifulSoup
from playwright.async_api import async_playwright, Page
from typing import Optional, List, Dict, Literal, Any
import random
from PIL import Image
import pytesseract
import shutil
import requests
import json
import fitz  # PyMuPDF for PDF
import tempfile
from pptx import Presentation
from prompt_manager import GoogleSheetPromptManager

# Load environment variables
load_dotenv(override=True)

# Pytesseract Path - 支援多種可能的路徑
tesseract_cmd = os.getenv('TESSERACT_CMD') or os.getenv('TESSERACT') 
tesseract_found = False

if tesseract_cmd:
    # 使用環境變數指定的路徑
    if os.path.exists(tesseract_cmd) or shutil.which(tesseract_cmd):
        pytesseract.pytesseract.tesseract_cmd = tesseract_cmd
        tesseract_found = True
    else:
        # 環境變數路徑無效，嘗試其他路徑
        tesseract_cmd = None

if not tesseract_found:
    # 嘗試常見路徑
    possible_paths = [
        '/root/.nix-profile/bin/tesseract',  # Railway/Nix 環境
        '/usr/bin/tesseract',  # Ubuntu/Debian
        '/usr/local/bin/tesseract',  # 標準 Unix 位置
        '/bin/tesseract',  # 某些 Linux 發行版
        '/opt/homebrew/bin/tesseract',  # macOS Homebrew
        'tesseract'  # PATH 中（最後嘗試）
    ]
    
    for path in possible_paths:
        if shutil.which(path) or os.path.exists(path):
            pytesseract.pytesseract.tesseract_cmd = path
            tesseract_found = True
            break
    
    if not tesseract_found:
        # 如果都找不到，使用預設值但會在後續處理中優雅降級
        pytesseract.pytesseract.tesseract_cmd = 'tesseract'

# 配置日誌
logger = logging.getLogger(__name__)
if tesseract_found:
    logger.info(f"✅ Tesseract OCR 找到: {pytesseract.pytesseract.tesseract_cmd}")
else:
    logger.warning(f"⚠️ Tesseract OCR 未找到，OCR 功能將被停用")

logger.setLevel(logging.INFO)

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s %(levelname)s %(message)s",
    handlers=[logging.StreamHandler()]
)

# 延遲初始化 prompt_manager
prompt_manager = None

class DeckBrowser:
    
    def __init__(self, prompt_manager: GoogleSheetPromptManager = None):
        """Initialize the DeckBrowser."""
        # 延遲初始化 prompt_manager，避免啟動時網路問題
        self.prompt_manager = prompt_manager
        
        # 設置日誌
        self.logger = logging.getLogger(__name__)
        self.browser = None
        self.email = os.getenv("DOCSEND_EMAIL")  # 替換為您的電子郵件
        self.path_helper = PathHelper()
        self.docsend_password = None  # 新增：儲存 DocSend 密碼

    #決定流程
    SourceType = Literal["docsend", "attachment", "gdrive", "website", "unknown"]
    
    @staticmethod
    def detect_source_type(message: str, attachments: Optional[list] = None) -> "DeckBrowser.SourceType":
        if "docsend.com" in message.lower():
            return "docsend"
        elif attachments and any(
            att.get("name", "").lower().endswith(('.pdf', '.pptx', '.ppt'))
            for att in attachments if isinstance(att, dict)
        ):
            return "attachment"
        elif re.search(r"https://(?:drive|docs)\.google\.com/(?:file/d/|presentation/)[\w\-/]+", message):
            return "gdrive"
        elif re.search(r"https?://[^\s\)]+", message):  # 匹配任何網址
            return "website"
        else:
            return "unknown"
    
    def extract_password_from_message(self, message: str) -> str:
        """從 message 中擷取 password: xxx 或 密碼: xxx"""
        import re
        patterns = [
            r'password[:：]?\s*([\w\-!@#$%^&*()_+=]+)',
            r'passwords[:：]?\s*([\w\-!@#$%^&*()_+=]+)',
            r'pw[:：]?\s*([\w\-!@#$%^&*()_+=]+)',
            r'密碼[:：]?\s*([\w\-!@#$%^&*()_+=]+)'
        ]
        for pattern in patterns:
            match = re.search(pattern, message, re.IGNORECASE)
            if match:
                return match.group(1)
        return None

    async def process_input(self, message: str, attachments: Optional[list] = None):
        # 先擷取密碼
        self.docsend_password = self.extract_password_from_message(message)
        self.logger.info(f"已擷取密碼: {self.docsend_password}" )
        results = []  # 用於存儲所有結果
        processed_urls = set()

        # 1. DocSend
        if "docsend.com" in message.lower():
            self.logger.info(f"開始處理 Docsend")
            docsend_results = await self.run_docsend_analysis(message)
            if docsend_results:
                results.extend(docsend_results)
            # 收集已處理過的 DocSend 連結
            docsend_urls = await self.extract_docsend_links(message)
            processed_urls.update(docsend_urls)

        # 2. Attachment
        if attachments and any(
            att.get("name", "").lower().endswith((".pdf", ".pptx", ".ppt"))
            for att in attachments if isinstance(att, dict)
        ):
            self.logger.info(f"開始處理 Attachment")
            attachment_results = await self.run_file_analysis(attachments)
            if attachment_results:
                results.extend(attachment_results)
            # 附件通常沒有網址，這裡略過

        # 3. GDrive
        gdrive_urls = []
        if re.search(r"https://(?:drive|docs)\.google\.com/(?:file/d/|presentation/)[\w\-/]+", message):
            self.logger.info(f"開始處理 Google Drive")
            gdrive_results = await self.run_gdrive_analysis(message)
            if gdrive_results:
                results.extend(gdrive_results)
            # 收集已處理過的 GDrive 連結
            gdrive_urls = re.findall(r'https://drive\.google\.com/file/d/[\w-]+|https://docs\.google\.com/presentation/d/[\w-]+', message)
            processed_urls.update(gdrive_urls)

        # 4. Generic Website
        if re.search(r"https?://[^\s\)]+", message):
            self.logger.info("🔗 偵測為一般網站，開始擷取網頁內容進行分析")
            generic_results = await self.run_generic_link_analysis(message, exclude_urls=processed_urls)
            if generic_results:
                results.extend(generic_results)

        # 5. 純文字
        if not results:
            self.logger.info("未偵測到連結或附件，直接分析純文字內容")
            summary = await summarize_pitch_deck(message, message)
            return [summary] if summary else [{"error": "❌ 純文字分析失敗"}]

        return results if results else [{"error": "❌ 沒有成功擷取任何內容"}]

    async def run_docsend_analysis(self, message: str) -> List[Dict[str, Any]]:
        """
        封裝完整流程：初始化 -> 擷取 URL -> 抽取內容 -> 關閉瀏覽器 -> 回傳結果
        """
        await self.initialize()
        
        results = []  # List to store JSON results
        urls = await self.extract_docsend_links(message)
        
        for url in urls:
            self.logger.info(f"處理 DocSend 連結: {url}")
            content = await self.read_docsend_document(url)
            if isinstance(content, dict):
                results.append(content)
            elif isinstance(content, str):
                summarized = await summarize_pitch_deck(content, message)
                if summarized:
                    results.append(summarized)

        await self.close()
        
        return results if results else [{"error": "❌ 沒有成功擷取任何 DocSend 文檔內容"}]

    async def run_gdrive_analysis(self, message: str) -> List[Dict[str, Any]]:
        self.logger.info(f"📥 開始處理 Google Drive 連結")
        results = []

        # 提取所有 Google Drive 連結
        gdrive_urls = re.findall(r'https://drive\.google\.com/file/d/([\w-]+)|id=([\w-]+)', message)
        
        for match in gdrive_urls:
            file_id = match[0] or match[1]  # 使用非空的匹配結果
            export_url = f"https://drive.google.com/uc?export=download&id={file_id}"
            
            try:
                response = requests.get(export_url)
                response.raise_for_status()

                with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
                    tmp.write(response.content)
                    tmp_path = tmp.name

                self.logger.info(f"📄 成功下載 PDF 檔案至 {tmp_path}，開始執行分析...")

                # 使用 Playwright 打開 PDF 文件
                await self.initialize()
                page = await self._get_page(f"file://{tmp_path}")
                
                try:
                    # 等待 PDF 加載
                    await page.wait_for_load_state('networkidle', timeout=30000)
                    
                    # 滾動頁面以加載所有內容
                    self.logger.info("開始滾動頁面以加載所有內容")
                    last_height = await page.evaluate('document.body.scrollHeight')
                    while True:
                        await page.evaluate('window.scrollTo(0, document.body.scrollHeight)')
                        await page.wait_for_timeout(2000)
                        new_height = await page.evaluate('document.body.scrollHeight')
                        if new_height == last_height:
                            break
                        last_height = new_height
                        self.logger.info("繼續滾動以加載更多內容")
                    
                    # 滾動回頂部
                    await page.evaluate('window.scrollTo(0, 0)')
                    await page.wait_for_timeout(1000)
                    
                    # 提取內容
                    content = await page.content()
                    soup = BeautifulSoup(content, 'html.parser')
                    
                    # 提取文本
                    text_elements = soup.find_all(['p', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'li', 'div'])
                    extracted_text = []
                    
                    for elem in text_elements:
                        text = elem.get_text().strip()
                        if text and len(text) > 10:
                            extracted_text.append(text)
                    
                    if not extracted_text:
                        # 如果沒有找到文本，嘗試 OCR
                        self.logger.warning("⚠️ 沒有找到文本內容，嘗試 OCR")
                        img_tags = soup.find_all('img')
                        image_urls = [img['src'] for img in img_tags if img.get('src')]
                        
                        if image_urls:
                            ocr_text = await ocr_images_from_urls(image_urls)
                            if ocr_text:
                                summary = await summarize_pitch_deck(ocr_text, message)
                                results.append(summary)
                                continue
                    
                    # 格式化內容
                    formatted_content = "\n\n".join(extracted_text)
                    summary = await summarize_pitch_deck(formatted_content, message)
                    if summary:
                        results.append(summary)
                    
                except Exception as e:
                    self.logger.error(f"處理 PDF 時出錯: {e}")
                    results.append({"error": f"❌ 處理 PDF 失敗: {str(e)}"})
                finally:
                    await page.close()
                    await self.close()
                
            except Exception as e:
                self.logger.error(f"❌ 下載或處理 Google Drive 檔案時發生錯誤：{str(e)}")
                results.append({"error": f"❌ Google Drive 檔案處理失敗：{str(e)}"})
        
        return results if results else [{"error": "❌ 沒有成功處理任何 Google Drive 檔案"}]

    async def run_file_analysis(self, attachments: List[Dict[str, str]]) -> List[Dict[str, Any]]:
        """
        處理 PDF/PPTX 附件：執行 OCR 或結構化摘要
        """
        results = []

        for file in attachments:
            path = file.get("path")
            name = file.get("name", "unnamed")
            suffix = self.path_helper.get(name).suffix.lower()

            self.logger.info(f"📂 開始分析附件: {name}")

            extracted_text = ""

            try:
                # Check file size before processing
                file_path_obj = self.path_helper.get(path)
                if not file_path_obj.exists() or file_path_obj.stat().st_size < 1024:
                    self.logger.error(f"❌ 檔案 {name} ({path}) 太小或不存在，可能下載失敗。")
                    results.append({"error": f"❌ 檔案 {name} 下載失敗或不是有效的檔案。"})
                    continue

                if suffix == ".pdf":
                    doc = fitz.open(str(file_path_obj))
                    for page in doc:
                        extracted_text += page.get_text("text") + "\n"
                    doc.close()

                elif suffix == ".pptx":
                    try:
                        prs = Presentation(str(file_path_obj))
                    except Exception as e:
                        self.logger.error(f"❌ 無法開啟 PPTX 檔案 {name}: {type(e).__name__}: {e}")
                        results.append({"error": f"❌ 無法開啟 PPTX 檔案 {name}: {type(e).__name__}: {e}"})
                        continue
                    for i, slide in enumerate(prs.slides):
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                extracted_text += f"[Slide {i+1}]\n{shape.text}\n"

                if not extracted_text.strip():
                    self.logger.warning(f"⚠️ {name} 沒有擷取到文字，開始執行 OCR 圖片辨識")

                    image_urls = []
                    if suffix == ".pdf":
                        doc = fitz.open(str(file_path_obj))
                        for page_num, page in enumerate(doc):
                            pix = page.get_pixmap(dpi=200)
                            img_path_obj = self.path_helper.get(f"{file_path_obj}_page_{page_num}.png")
                            pix.save(str(img_path_obj))
                            image_urls.append(f"file://{img_path_obj}")
                        doc.close()

                    elif suffix == ".pptx":
                        self.logger.warning("❌ 尚未實作 PPTX 頁面轉圖片的 OCR fallback")
                        continue

                    if image_urls:
                        ocr_text = await ocr_images_from_urls(image_urls)
                        if ocr_text:
                            summary = await summarize_pitch_deck(ocr_text, name)
                            if summary:
                                results.append(summary)
                                continue

                else:
                    summary = await summarize_pitch_deck(extracted_text, name)
                    if summary:
                        results.append(summary)

            except Exception as e:
                self.logger.error(f"❌ 分析檔案 {name} 發生錯誤：{e}")
                results.append({"error": f"❌ 分析失敗: {name}"})

        return results if results else [{"error": "❌ 沒有成功處理任何附件內容"}]

    async def initialize(self):
        """Initialize the browser and browser context asynchronously."""
        try:
            self.playwright = await async_playwright().start()
            self.browser = await self.playwright.chromium.launch(
                headless=True,
                args=[
                    "--disable-gpu",
                    "--no-sandbox",
                    "--disable-dev-shm-usage",
                    "--disable-setuid-sandbox",
                    "--disable-extensions",
                    "--disable-background-networking",
                    "--disable-sync",
                    "--metrics-recording-only",
                    "--disable-default-apps",
                    "--mute-audio",
                    "--no-first-run",
                    "--hide-scrollbars",
                    "--ignore-certificate-errors",
                    "--window-size=1280,800",
                    "--single-process",
                    "--disable-blink-features=AutomationControlled"
                ]
            )

            # ✅ 建議在這裡建立 context，而不是延後到 _get_page 才建
            self.context = await self.browser.new_context(
                user_agent="Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/123.0.0.0 Safari/537.36",
                viewport={"width": 1280, "height": 800},
                locale="en-US",
                timezone_id="America/Los_Angeles",
                permissions=["geolocation"],
                extra_http_headers={
                    "Accept-Language": "en-US,en;q=0.9",
                    "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8",
                    "sec-ch-ua": '"Chromium";v="123", "Not:A-Brand";v="99", "Google Chrome";v="123"',
                    "sec-ch-ua-mobile": "?0",
                    "sec-ch-ua-platform": '"Windows"',
                }
            )

            self.logger.info("✅ Browser + context initialized")
        except Exception as e:
            self.logger.error(f"❌ Failed to initialize browser: {str(e)}")
            raise



    async def close(self):
        """關閉瀏覽器"""
        try:
            if self.browser:
                await self.browser.close()
                self.browser = None
            if hasattr(self, 'playwright'):
                await self.playwright.stop()
        except Exception as e:
            self.logger.error(f"Error closing browser: {e}")
        finally:
            # 確保所有資源都被釋放
            import gc
            gc.collect()

    async def extract_docsend_links(self, text: str) -> List[str]:
        """從文本中提取 DocSend 連結"""
        docsend_pattern = r'https?://(?:www\.)?docsend\.com/[^\s)"}]+'
        return re.findall(docsend_pattern, text)
    
    async def _get_page(self, url: str) -> Page:
        if not self.context:
            raise RuntimeError("Browser context not initialized")
        page = await self.context.new_page()
        page.set_default_timeout(30000)
        return page


    async def read_docsend_document(self, url: str) -> Optional[str]:
        """讀取 DocSend 文檔的內容"""
        try:
            self.logger.info(f"正在從 DocSend 連結讀取內容: {url}")
            
            # 創建新頁面
            page = await self._get_page(url)
            
            # 訪問 DocSend 頁面
            response = await page.goto(url, wait_until='networkidle', timeout=30000)
            
            self.logger.info(f"訪問狀態碼: {response.status}")
            
            if response.status == 403:
                logger.error("訪問被拒絕 (403 Forbidden)")
                await page.close()
                return None
            
            # 隨機等待 2-5 秒
            await asyncio.sleep(random.uniform(2, 5))
            
            # 檢查是否需要填寫電子郵件
            self.logger.info("[DocSend] 準備檢查是否需要填寫 email")
            email_input = await page.query_selector('input[type="email"]')
            if email_input:
                self.logger.info("[DocSend] 偵測到 email 輸入框，準備填寫 email")
                await page.type('input[type="email"]', self.email, delay=random.uniform(100, 200))
                self.logger.info(f"[DocSend] 已輸入 email: {self.email}")
                # 填完 email 後找密碼 input，支援多種 selector
                selectors = [
                    'input.js-auth-form_passcode',
                    'input[name="link_auth_form[passcode]"]',
                    'input[id="link_auth_form_passcode"]',
                    'input[type="password"]'
                ]
                password_input = None
                used_selector = None
                for sel in selectors:
                    try:
                        password_input = await page.query_selector(sel)
                        if password_input:
                            self.logger.info(f"[DocSend] 用 selector {sel} 找到密碼欄位")
                            used_selector = sel
                            break
                    except Exception as e:
                        self.logger.warning(f"[DocSend] 用 selector {sel} 找密碼欄位失敗: {e}")
                if password_input and self.docsend_password:
                    await page.type(used_selector, self.docsend_password, delay=random.uniform(100, 200))
                    self.logger.info(f"[DocSend] 已用 {used_selector} 填入密碼: {self.docsend_password}")
                else:
                    self.logger.warning("[DocSend] 沒有找到可填寫的密碼欄位")
                # 填完密碼再按 Continue
                try:
                    btn = page.locator('button:has-text("Continue")')
                    await btn.click()
                    self.logger.info("[DocSend] 已點擊 Continue (email+password)")
                    try:
                        await page.wait_for_load_state('networkidle', timeout=8000)
                    except Exception as e:
                        self.logger.warning(f"[DocSend] 點擊 Continue 後等待 networkidle 超時: {e}，改用 sleep 3 秒")
                        await page.wait_for_timeout(3000)
                except Exception as e:
                    self.logger.warning(f"[DocSend] 提交 email+password 按鈕點擊失敗: {e}")
                self.logger.info("[DocSend] email+password 流程結束，進入下一步")
            else:
                # 沒有 password input，才按 Continue
                try:
                    self.logger.info("[DocSend] 沒有偵測到 password 輸入框，準備點擊 Continue 按鈕 (email only)")
                    await page.locator('button:has-text("Continue")').wait_for(state='visible', timeout=1000)
                    await page.locator('button:has-text("Continue")').click(timeout=1000)
                    self.logger.info("[DocSend] 已點擊 Continue (email only)")
                except Exception as e:
                    self.logger.warning(f"[DocSend] 提交 email 按鈕點擊失敗: {e}")
                await page.wait_for_load_state('networkidle', timeout=10000)
                self.logger.info("[DocSend] email only 流程結束，進入下一步")
            # --- 密碼自動填寫結束 ---

            # 滾動頁面以加載所有內容
            self.logger.info("開始滾動頁面以加載所有內容")
            last_height = await page.evaluate('document.body.scrollHeight')
            while True:
                # 滾動到底部
                await page.evaluate('window.scrollTo(0, document.body.scrollHeight)')
                # 等待新內容加載
                await page.wait_for_timeout(2000)
                # 計算新的滾動高度
                new_height = await page.evaluate('document.body.scrollHeight')
                if new_height == last_height:
                    break
                last_height = new_height
                self.logger.info("繼續滾動以加載更多內容")
            
            # 滾動回頂部
            await page.evaluate('window.scrollTo(0, 0)')
            await page.wait_for_timeout(1000)
            
            # 嘗試點擊下一頁按鈕加載更多內容
            self.logger.info("嘗試通過點擊下一頁按鈕加載更多內容")
            try:
                next_buttons = await page.query_selector_all('button[aria-label*="next"], button[aria-label*="Next"], button[title*="next"], button[title*="Next"]')
                for button in next_buttons:
                    for _ in range(10):  # 最多點擊10次
                        try:
                            await button.click()
                            await page.wait_for_timeout(1000)  # 等待新頁面加載
                        except:
                            break
            except Exception as e:
                self.logger.warning(f"點擊下一頁按鈕時出錯: {e}")
            
            # 再次滾動以確保所有內容都已加載
            await page.evaluate('window.scrollTo(0, document.body.scrollHeight)')
            await page.wait_for_timeout(2000)
            
            # 截圖以便調試
            debug_screenshot_name = f"docsend_debug_{random.randint(1000, 9999)}.png"
            debug_screenshot_path = str(self.path_helper.get("tmp", debug_screenshot_name))
            self.path_helper.ensure_dir("tmp")
            await page.screenshot(path=debug_screenshot_path)
            self.logger.info(f"保存頁面截圖至: {debug_screenshot_path}")
            
            # Debug 將整頁 HTML 儲存下來
            html = await page.content()
            debug_html_path = self.path_helper.get("debug_docsend.html")
            with debug_html_path.open("w", encoding="utf-8") as f:
                f.write(html)

            # 等待文檔內容加載
            target_frame = None
            content_selectors = ['.document-content', '.page', '.viewer-content', '.ds-viewer-container']

            for frame in page.frames:
                frame_url = frame.url or ""
                if "docsend.com/view" in frame_url and "marketing.docsend.com" not in frame_url:
                    for selector in content_selectors:
                        try:
                            await frame.wait_for_selector(selector, timeout=5000)
                            self.logger.info(f"在 iframe 中找到內容: {selector}")
                            target_frame = frame
                            break
                        except:
                            continue
                if target_frame:
                    break

            if not target_frame:
                self.logger.warning("⚠️ 未能找到含文字 selector，開始 debug 所有 iframe...")
                await debug_all_iframes(page)

                # 嘗試從所有 iframe 中擷取圖片做 OCR
                for frame in page.frames:
                    try:
                        frame_html = await frame.content()
                        soup = BeautifulSoup(frame_html, "html.parser")
                        img_tags = soup.find_all("img")
                        image_urls = [img["src"] for img in img_tags if img.get("src")]
                        
                        if image_urls:
                            self.logger.info(f"✅ 在 iframe 中找到圖片：共 {len(image_urls)} 張，準備執行 OCR")
                            ocr_raw_text = await ocr_images_from_urls(image_urls)
                            if ocr_raw_text:
                                summarized_text = await summarize_pitch_deck(ocr_raw_text, url)
                                await page.close()
                                return summarized_text
                    except Exception as e:
                        self.logger.warning(f"讀取 iframe 內容失敗：{e}")

                self.logger.error("❌ 所有 iframe 均未能提供內容或圖片")
                await page.close()
                return None
            
            # 從 iframe 抓 HTML
            html = await target_frame.content()
            self.logger.info(f"獲取的HTML內容長度: {len(html)} 字符")
            
            # 解析 HTML
            soup = BeautifulSoup(html, 'html.parser')
            
            # 提取文檔標題
            title_elem = soup.find('title')
            if not title_elem:
                self.logger.info("iframe 中未找到 <title>，嘗試從主頁抓取 title")
                main_html = await page.content()
                main_soup = BeautifulSoup(main_html, 'html.parser')
                title_elem = main_soup.find('title')

            title = title_elem.text.strip() if title_elem else "DocSend Document"
            self.logger.info(f"提取的文檔標題: {title}")
            
            # 提取文檔內容
            text_elements = soup.find_all(['p', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'li', 'div'])
            self.logger.info(f"找到 {len(text_elements)} 個文本元素")
            
            extracted_text = []
            
            for elem in text_elements:
                text = elem.get_text().strip()
                if text and len(text) > 10:  # 只保留有意義的文本
                    extracted_text.append(text)
            
            self.logger.info(f"篩選後提取了 {len(extracted_text)} 個有意義的文本段落")
            extracted_text = "\n\n".join(extracted_text)
            
            # 關閉頁面
            await page.close()
            
            if not extracted_text or len(extracted_text) < 100:
                self.logger.warning("⚠️ 文字內容過少，嘗試 OCR 圖片並用 GPT 摘要")
                img_tags = soup.find_all('img')
                image_urls = [img['src'] for img in img_tags if img.get('src')]

                if image_urls:
                    ocr_raw_text = await ocr_images_from_urls(image_urls)
                    if ocr_raw_text:
                        summarized_text = await summarize_pitch_deck(ocr_raw_text, url)
                        extracted_text = summarized_text
            
            # 記錄內容摘要
            content_preview = extracted_text[:100] + "..." if len(extracted_text) > 100 else extracted_text
            self.logger.info(f"提取文本預覽: {content_preview}")
            self.logger.info(f"總提取文本長度: {len(extracted_text)} 字符")
            
            # 格式化提取的內容
            formatted_content = f"--- DocSend 文檔: {title} ---\n\n{extracted_text}\n\n--- DocSend 文檔結束 ---"
            
            self.logger.info(f"成功從 DocSend 連結提取內容: {url}")
            return formatted_content
        
        except Exception as e:
            self.logger.error(f"讀取 DocSend 文檔時出錯: {str(e)}", exc_info=True)
            if 'page' in locals():
                await page.close()
            return None

    async def is_pitch_deck(self, page) -> bool:
        """判斷網頁是否為 Pitch Deck"""
        try:
            # 0. 首先檢查是否為社群網站
            url = page.url.lower()
            social_media_domains = ['x.com', 'twitter.com', 'facebook.com', 'linkedin.com', 'instagram.com']
            if any(domain in url for domain in social_media_domains):
                return False

            # 1. 檢查 URL 關鍵字
            deck_keywords = ['pitch', 'deck', 'presentation', 'slides', 'investor', 'fundraising']
            if any(keyword in url for keyword in deck_keywords):
                return True

            # 2. 檢查頁面標題
            title = await page.title()
            if title and any(keyword in title.lower() for keyword in deck_keywords):
                return True

            # 3. 檢查頁面內容特徵
            content = await page.content()
            soup = BeautifulSoup(content, "html.parser")
            
            # 檢查是否有投影片相關的元素
            slide_indicators = [
                # 檢查是否有投影片容器
                bool(soup.find('div', class_=lambda x: x and any(word in str(x).lower() for word in ['slide', 'deck', 'presentation']))),
                # 檢查是否有投影片導航
                bool(soup.find('nav', class_=lambda x: x and any(word in str(x).lower() for word in ['slide', 'deck', 'presentation']))),
                # 檢查是否有投影片計數器
                bool(soup.find('div', class_=lambda x: x and any(word in str(x).lower() for word in ['counter', 'progress', 'slide-number']))),
                # 檢查是否有投影片控制按鈕
                bool(soup.find('button', class_=lambda x: x and any(word in str(x).lower() for word in ['next', 'prev', 'slide']))),
            ]
            
            if any(slide_indicators):
                return True

            # 4. 檢查頁面結構
            # 計算頁面中的圖片數量
            images = await page.query_selector_all('img')
            # 計算頁面中的文本塊數量
            text_blocks = await page.query_selector_all('p, h1, h2, h3, h4, h5, h6')
            
            # 如果圖片數量較多且文本塊較少，可能是投影片
            # 但需要確保不是社群網站的結構
            if len(images) > 5 and len(text_blocks) < 10:
                # 檢查是否有社群網站特有的元素
                social_indicators = [
                    bool(soup.find('div', class_=lambda x: x and any(word in str(x).lower() for word in ['tweet', 'post', 'status', 'feed', 'timeline']))),
                    bool(soup.find('div', class_=lambda x: x and any(word in str(x).lower() for word in ['profile', 'avatar', 'user-info']))),
                    bool(soup.find('div', class_=lambda x: x and any(word in str(x).lower() for word in ['like', 'share', 'comment', 'retweet']))),
                ]
                if not any(social_indicators):
                    return True

            # 5. 檢查是否有投影片相關的 JavaScript
            scripts = await page.query_selector_all('script')
            for script in scripts:
                content = await script.text_content()
                if any(keyword in content.lower() for keyword in ['slideshow', 'presentation', 'deck']):
                    return True

            return False

        except Exception as e:
            self.logger.error(f"Error checking if page is pitch deck: {e}")
            return False

    async def run_generic_link_analysis(self, message: str, exclude_urls: set = None) -> List[Dict[str, Any]]:
        """分析一般網址（包括公司官網），每個網址單獨回傳 summary，不合併統整"""
        urls = re.findall(r'https?://[^\s\)]+', message)
        if exclude_urls:
            # 標準化網址（去除末尾斜線）
            def norm(u):
                return u.rstrip('/')
            exclude_set = set(map(norm, exclude_urls))
            urls = [u for u in urls if norm(u) not in exclude_set]
        results = []
        
        if not urls:
            return [{"error": "❌ 未找到任何有效的網址"}]
        
        self.logger.info(f"找到 {len(urls)} 個網址需要處理")
        
        for url in urls:
            try:
                self.logger.info(f"🌐 開始分析網址: {url}")
                content = await self.extract_content(url)
                if content:
                    self.logger.info(f"成功提取內容: {len(content)} 字符")
                    results.append({"url": url, "summary": content})
                else:
                    self.logger.warning("❌ 沒有提取到內容")
                    results.append({"url": url, "error": "❌ 無法提取內容"})
            except Exception as e:
                self.logger.error(f"❌ 分析 {url} 失敗：{e}")
                results.append({"url": url, "error": f"❌ 分析失敗: {e}"})
        return results if results else [{"error": "❌ 沒有成功處理任何網址"}]

    async def process_pitch_deck_page(self, page) -> Optional[str]:
        """處理 Pitch Deck 頁面"""
        try:
            self.logger.info("開始處理 Pitch Deck 頁面")
            
            # 檢查是否為 Journey.io 網站
            if "journey.io" in page.url:
                self.logger.info("檢測到 Journey.io 網站，使用特殊處理方式")
                return await self.process_journey_page(page)
            
            # 等待頁面加載完成
            await page.wait_for_load_state('networkidle', timeout=30000)
            
            # 截圖以便調試
            debug_screenshot_name = f"pitch_deck_debug_{random.randint(1000, 9999)}.png"
            debug_screenshot_path = str(self.path_helper.get("tmp", debug_screenshot_name))
            self.path_helper.ensure_dir("tmp")
            await page.screenshot(path=debug_screenshot_path)
            self.logger.info(f"保存頁面截圖至: {debug_screenshot_path}")
            
            # Debug 將整頁 HTML 儲存下來
            html = await page.content()
            debug_html_path = self.path_helper.get("debug_pitch_deck.html")
            with debug_html_path.open("w", encoding="utf-8") as f:
                f.write(html)
            
            # 提取頁面標題
            title = await page.title()
            self.logger.info(f"提取的頁面標題: {title}")
            
            # 1. 嘗試提取投影片內容
            slides = []
            
            # 等待投影片容器加載
            try:
                # 擴展選擇器以匹配更多可能的投影片容器
                await page.wait_for_selector('div[class*="slide"], div[class*="deck"], div[class*="presentation"], div[class*="page"], div[class*="slide-container"], div[class*="slide-wrapper"]', timeout=5000)
            except Exception as e:
                self.logger.warning(f"等待投影片容器超時: {e}")
            
            # 滾動頁面以加載所有投影片
            self.logger.info("開始滾動頁面以加載所有投影片")
            last_height = await page.evaluate('document.body.scrollHeight')
            while True:
                # 滾動到底部
                await page.evaluate('window.scrollTo(0, document.body.scrollHeight)')
                # 等待新內容加載
                await page.wait_for_timeout(2000)
                # 計算新的滾動高度
                new_height = await page.evaluate('document.body.scrollHeight')
                if new_height == last_height:
                    break
                last_height = new_height
                self.logger.info("繼續滾動以加載更多投影片")
            
            # 滾動回頂部
            await page.evaluate('window.scrollTo(0, 0)')
            await page.wait_for_timeout(1000)
            
            # 獲取所有投影片
            slide_elements = await page.query_selector_all('div[class*="slide"], div[class*="deck"], div[class*="presentation"], div[class*="page"], div[class*="slide-container"], div[class*="slide-wrapper"]')
            self.logger.info(f"找到 {len(slide_elements)} 個投影片元素")
            
            # 如果找不到投影片，嘗試其他選擇器
            if not slide_elements:
                self.logger.info("嘗試使用備用選擇器")
                # 嘗試查找所有可能的投影片容器
                slide_elements = await page.query_selector_all('div[role="slide"], div[data-slide], div[data-page], div[data-index]')
                self.logger.info(f"使用備用選擇器找到 {len(slide_elements)} 個投影片元素")
            
            # 如果還是找不到，嘗試查找所有可能包含投影片內容的 div
            if not slide_elements:
                self.logger.info("嘗試查找所有可能包含投影片內容的 div")
                # 獲取頁面中的所有 div
                all_divs = await page.query_selector_all('div')
                # 過濾出可能包含投影片內容的 div
                slide_elements = []
                for div in all_divs:
                    try:
                        # 檢查 div 是否包含投影片相關的內容
                        text = await div.text_content()
                        if text and len(text.strip()) > 50:  # 假設投影片內容通常較長
                            slide_elements.append(div)
                    except:
                        continue
                self.logger.info(f"通過內容分析找到 {len(slide_elements)} 個可能的投影片")
            
            # 如果找到的投影片數量明顯少於預期，嘗試點擊下一頁按鈕
            if len(slide_elements) < 10:  # 假設至少有10頁
                self.logger.info("嘗試通過點擊下一頁按鈕加載更多投影片")
                try:
                    # 嘗試點擊下一頁按鈕
                    next_buttons = await page.query_selector_all('button[aria-label*="next"], button[aria-label*="Next"], button[title*="next"], button[title*="Next"]')
                    for button in next_buttons:
                        for _ in range(10):  # 最多點擊10次
                            try:
                                await button.click()
                                await page.wait_for_timeout(1000)  # 等待新頁面加載
                            except:
                                break
                except Exception as e:
                    self.logger.warning(f"點擊下一頁按鈕時出錯: {e}")
                
                # 重新獲取投影片
                slide_elements = await page.query_selector_all('div[class*="slide"], div[class*="deck"], div[class*="presentation"], div[class*="page"], div[class*="slide-container"], div[class*="slide-wrapper"]')
                self.logger.info(f"點擊下一頁後找到 {len(slide_elements)} 個投影片元素")
            
            for i, slide in enumerate(slide_elements):
                try:
                    # 提取投影片文本
                    text = await slide.text_content()
                    if text.strip():
                        slides.append(f"[Slide {i+1}]\n{text.strip()}")
                        self.logger.info(f"成功提取第 {i+1} 張投影片的文本")
                    
                    # 提取投影片圖片
                    images = await slide.query_selector_all('img')
                    self.logger.info(f"第 {i+1} 張投影片中找到 {len(images)} 張圖片")
                    
                    for img in images:
                        src = await img.get_attribute('src')
                        if src:
                            try:
                                # 下載圖片並進行 OCR
                                response = requests.get(src)
                                img = Image.open(BytesIO(response.content))
                                
                                try:
                                    ocr_text = pytesseract.image_to_string(img)
                                    if ocr_text.strip():
                                        slides.append(f"[Slide {i+1} Image]\n{ocr_text.strip()}")
                                        self.logger.info(f"成功對第 {i+1} 張投影片的圖片進行 OCR")
                                except Exception as ocr_e:
                                    # 處理 Tesseract 相關錯誤
                                    if "TesseractNotFoundError" in str(type(ocr_e)) or "tesseract" in str(ocr_e).lower():
                                        self.logger.warning(f"⚠️ Tesseract OCR 不可用，跳過第 {i+1} 張投影片圖片的文字提取")
                                        slides.append(f"[Slide {i+1} Image]\n[OCR不可用 - 無法提取文字內容]")
                                    else:
                                        self.logger.warning(f"❌ OCR 處理失敗 (第 {i+1} 張): {str(ocr_e)}")
                                        slides.append(f"[Slide {i+1} Image]\n[文字提取失敗]")
                                
                            except Exception as e:
                                self.logger.warning(f"OCR 失敗: {e}")
                except Exception as e:
                    self.logger.warning(f"處理第 {i+1} 張投影片時出錯: {e}")
            
            if slides:
                formatted_content = f"--- Pitch Deck: {title} ---\n\n" + "\n\n".join(slides) + "\n\n--- Pitch Deck 結束 ---"
                self.logger.info("成功提取投影片內容")
                return formatted_content
            
            # 2. 如果無法提取投影片，嘗試提取整個頁面內容
            self.logger.info("無法提取投影片內容，嘗試提取整個頁面")
            content = await page.content()
            soup = BeautifulSoup(content, "html.parser")
            
            # 提取所有文本內容
            texts = []
            for element in soup.find_all(['p', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'li', 'div']):
                text = element.get_text(strip=True)
                if text and len(text) > 10:
                    texts.append(text)
            
            if texts:
                formatted_content = f"--- Pitch Deck: {title} ---\n\n" + "\n\n".join(texts) + "\n\n--- Pitch Deck 結束 ---"
                self.logger.info("成功提取頁面文本內容")
                return formatted_content
            
            # 3. 如果文字提取失敗，嘗試 OCR 所有圖片
            self.logger.warning("文字提取失敗，嘗試 OCR 所有圖片")
            img_tags = soup.find_all("img")
            image_urls = [img.get("src") for img in img_tags if img.get("src")]
            
            if image_urls:
                self.logger.info(f"找到 {len(image_urls)} 張圖片，開始 OCR")
                ocr_text = await ocr_images_from_urls(image_urls)
                if ocr_text.strip():
                    formatted_content = f"--- Pitch Deck: {title} ---\n\n{ocr_text}\n\n--- Pitch Deck 結束 ---"
                    self.logger.info("成功通過 OCR 提取內容")
                    return formatted_content
            
            self.logger.error("❌ 無法提取任何有效內容")
            return None
            
        except Exception as e:
            self.logger.error(f"處理 Pitch Deck 頁面時出錯: {str(e)}", exc_info=True)
            return None

    async def process_journey_page(self, page) -> Optional[str]:
        """處理 Journey.io 頁面"""
        try:
            self.logger.info("開始處理 Journey.io 頁面")
            
            # 等待頁面加載
            await page.wait_for_load_state('networkidle', timeout=30000)
            
            # 保存頁面結構以供調試
            debug_html_path = self.path_helper.get("debug_journey.html")
            html_content = await page.content()
            with debug_html_path.open("w", encoding="utf-8") as f:
                f.write(html_content)
            self.logger.info(f"已保存頁面結構到 {debug_html_path}")
            
            # 保存頁面截圖
            debug_screenshot_path = self.path_helper.get("debug_journey.png")
            await page.screenshot(path=str(debug_screenshot_path))
            self.logger.info(f"已保存頁面截圖到 {debug_screenshot_path}")
            
            # 提取頁面標題
            title = await page.title()
            self.logger.info(f"提取的頁面標題: {title}")
            
            # 改進的滾動邏輯
            self.logger.info("開始滾動頁面以加載所有內容")
            
            # 1. 先滾動到頂部
            await page.evaluate('window.scrollTo(0, 0)')
            await page.wait_for_timeout(1000)
            
            # 2. 獲取初始頁面高度
            initial_height = await page.evaluate('document.body.scrollHeight')
            self.logger.info(f"初始頁面高度: {initial_height}")
            
            # 3. 使用更小的滾動步長和更長的等待時間
            scroll_step = 300  # 每次滾動300像素
            max_attempts = 50  # 最大嘗試次數
            no_change_count = 0  # 記錄高度未變化的次數
            
            for attempt in range(max_attempts):
                # 獲取當前滾動位置
                current_position = await page.evaluate('window.pageYOffset')
                
                # 計算新的滾動位置
                new_position = current_position + scroll_step
                
                # 滾動到新位置
                await page.evaluate(f'window.scrollTo(0, {new_position})')
                await page.wait_for_timeout(2000)  # 等待2秒讓內容加載
                
                # 獲取新的頁面高度
                new_height = await page.evaluate('document.body.scrollHeight')
                
                # 檢查是否到達底部
                if new_position >= new_height:
                    self.logger.info("已到達頁面底部")
                    break
                
                # 檢查頁面高度是否變化
                if new_height == initial_height:
                    no_change_count += 1
                    if no_change_count >= 3:  # 如果連續3次高度未變化，可能已經到底
                        self.logger.info("頁面高度連續未變化，可能已到底部")
                        break
                else:
                    no_change_count = 0
                    initial_height = new_height
                
                self.logger.info(f"滾動進度: {new_position}/{new_height} (嘗試 {attempt + 1}/{max_attempts})")
            
            # 4. 最後再滾動回頂部
            await page.evaluate('window.scrollTo(0, 0)')
            await page.wait_for_timeout(1000)
            
            # 使用 JavaScript 提取所有內容
            content = await page.evaluate('''() => {
                const sections = [];
                
                // 遍歷所有主要區塊
                document.querySelectorAll('div[class*="page"], div[class*="section"], div[class*="content"], div[class*="block"], div[class*="slide"]').forEach(section => {
                    const sectionData = {
                        title: '',
                        content: [],
                        teamMembers: []
                    };
                    
                    // 提取標題
                    const titleElem = section.querySelector('h1, h2, h3, h4, h5, h6, [class*="title"], [class*="heading"], [class*="header"]');
                    if (titleElem) {
                        sectionData.title = titleElem.textContent.trim();
                    }
                    
                    // 檢查是否為團隊相關部分
                    const isTeamSection = sectionData.title.toLowerCase().includes('team') || 
                                        sectionData.title.toLowerCase().includes('about us') ||
                                        sectionData.title.toLowerCase().includes('founder') ||
                                        sectionData.title.toLowerCase().includes('leadership');
                    
                    if (isTeamSection) {
                        console.log('Found team section:', sectionData.title);
                        
                        // 擴展團隊成員選擇器
                        const teamElements = section.querySelectorAll(`
                            div[class*="member"], 
                            div[class*="person"], 
                            div[class*="profile"], 
                            div[class*="team"], 
                            div[class*="founder"], 
                            div[class*="investor"], 
                            div[class*="funding"],
                            div[class*="card"],
                            div[class*="bio"],
                            div[class*="staff"],
                            div[class*="leadership"],
                            div[class*="executive"]
                        `);
                        
                        console.log('Found team elements:', teamElements.length);
                        
                        teamElements.forEach(teamElem => {
                            const memberData = {
                                name: '',
                                role: '',
                                description: ''
                            };
                            
                            // 擴展名稱選擇器
                            const nameElem = teamElem.querySelector(`
                                h1, h2, h3, h4, h5, h6, 
                                [class*="name"], 
                                [class*="title"],
                                [class*="heading"],
                                [class*="header"],
                                strong,
                                b
                            `);
                            if (nameElem) {
                                memberData.name = nameElem.textContent.trim();
                                console.log('Found team member name:', memberData.name);
                            }
                            
                            // 擴展角色選擇器
                            const roleElem = teamElem.querySelector(`
                                [class*="role"], 
                                [class*="position"], 
                                [class*="title"],
                                [class*="job"],
                                [class*="position"],
                                [class*="designation"],
                                em,
                                i
                            `);
                            if (roleElem) {
                                memberData.role = roleElem.textContent.trim();
                                console.log('Found team member role:', memberData.role);
                            }
                            
                            // 擴展描述選擇器
                            const descElem = teamElem.querySelector(`
                                p, 
                                div[class*="text"], 
                                div[class*="content"], 
                                div[class*="description"],
                                div[class*="bio"],
                                div[class*="about"],
                                div[class*="info"]
                            `);
                            if (descElem) {
                                memberData.description = descElem.textContent.trim();
                                console.log('Found team member description:', memberData.description);
                            }
                            
                            // 如果找到任何信息，就添加到團隊成員列表
                            if (memberData.name || memberData.role || memberData.description) {
                                sectionData.teamMembers.push(memberData);
                                console.log('Added team member:', memberData);
                            }
                        });
                    }
                    
                    // 提取所有文本內容
                    const textElements = section.querySelectorAll(`
                        p, 
                        li, 
                        div[class*="text"], 
                        div[class*="content"], 
                        div[class*="description"],
                        div[class*="info"],
                        div[class*="detail"]
                    `);
                    textElements.forEach(elem => {
                        const text = elem.textContent.trim();
                        if (text) {
                            sectionData.content.push(text);
                        }
                    });
                    
                    if (sectionData.title || sectionData.content.length > 0 || sectionData.teamMembers.length > 0) {
                        sections.push(sectionData);
                    }
                });
                
                return sections;
            }''')
            
            # 格式化提取的內容
            formatted_sections = []
            for section in content:
                section_text = []
                
                if section['title']:
                    section_text.append(f"\n## {section['title']}")
                    self.logger.info(f"Found section: {section['title']}")
                
                if section['content']:
                    section_text.extend(section['content'])
                    self.logger.info(f"Found {len(section['content'])} content items")
                
                if section['teamMembers']:
                    section_text.append("\n### Team Members:")
                    self.logger.info(f"Found {len(section['teamMembers'])} team members")
                    for member in section['teamMembers']:
                        member_text = []
                        if member['name']:
                            member_text.append(f"Name: {member['name']}")
                            self.logger.info(f"Found team member: {member['name']}")
                        if member['role']:
                            member_text.append(f"Role: {member['role']}")
                            self.logger.info(f"Member role: {member['role']}")
                        if member['description']:
                            member_text.append(f"Description: {member['description']}")
                            self.logger.info(f"Member description: {member['description']}")
                        if member_text:
                            section_text.append("\n" + "\n".join(member_text))
                
                if section_text:
                    formatted_sections.append('\n'.join(section_text))
            
            if formatted_sections:
                formatted_content = f"--- Journey.io Document: {title} ---\n\n" + "\n\n".join(formatted_sections) + "\n\n--- Document End ---"
                self.logger.info("成功提取 Journey.io 頁面內容")
                return formatted_content
            
            self.logger.error("❌ 無法提取任何有效內容")
            return None
            
        except Exception as e:
            self.logger.error(f"處理 Journey.io 頁面時出錯: {str(e)}", exc_info=True)
            return None

    async def extract_content(self, url: str) -> Optional[str]:
        """用 Playwright 取得渲染後內容，只抓主要文字內容，不呼叫 GPT"""
        try:
            self.logger.info(f"用 Playwright 取得渲染後內容: {url}")
            async with async_playwright() as pw:
                browser = await pw.chromium.launch(headless=True)
                page = await browser.new_page()
                try:
                    await page.goto(url, wait_until="domcontentloaded", timeout=60000)
                except Exception as e:
                    self.logger.warning(f"第一次 goto domcontentloaded 失敗: {e}")
                    try:
                        await page.goto(url, wait_until="load", timeout=60000)
                    except Exception as e2:
                        self.logger.error(f"第二次 goto load 也失敗: {e2}")
                        await browser.close()
                        return None

                # 判斷是否為 GitBook/Notion/Docs 這類有目錄的網站
                is_multi_page = False
                sidebar_selectors = [
                    'nav.toc a',           # GitBook
                    'nav[aria-label="Table of contents"] a',
                    'aside a',             # Notion/Docs
                    '.sidebar a',
                    '.menu a',
                    '.toc a',
                    'nav a',
                ]
                all_links = set()
                for sel in sidebar_selectors:
                    try:
                        links = await page.query_selector_all(sel)
                        for link in links:
                            href = await link.get_attribute('href')
                            if href and not href.startswith('#') and not href.startswith('javascript:'):
                                # 統一補全相對路徑
                                if href.startswith('/') and url.startswith('http'):
                                    from urllib.parse import urljoin
                                    href = urljoin(url, href)
                                elif href.startswith('http'):
                                    pass
                                else:
                                    continue
                                all_links.add(href)
                        if len(all_links) > 3:
                            is_multi_page = True
                            self.logger.info(f"偵測到多分頁目錄 selector: {sel}，共 {len(all_links)} 個分頁")
                            break
                    except Exception:
                        continue

                all_links = list(all_links)
                if url not in all_links:
                    all_links = [url] + all_links

                all_contents = []
                if is_multi_page:
                    # 遍歷所有分頁
                    for idx, link in enumerate(all_links):
                        try:
                            self.logger.info(f"[多分頁] 抓取第{idx+1}/{len(all_links)}頁: {link}")
                            await page.goto(link, wait_until="domcontentloaded", timeout=30000)
                            await page.wait_for_timeout(1000)
                            # 滾動到底部
                            await page.evaluate('window.scrollTo(0, document.body.scrollHeight)')
                            await page.wait_for_timeout(1000)
                            await page.evaluate('window.scrollTo(0, 0)')
                            await page.wait_for_timeout(500)
                            html = await page.content()
                            soup = BeautifulSoup(html, "html.parser")
                            title = soup.title.string.strip() if soup.title else ""
                            main = soup.find("main") or soup.find("article") or soup.body
                            text_blocks = []
                            if main:
                                for tag in main.find_all(["h1", "h2", "h3", "h4", "h5", "h6", "p", "li"]):
                                    txt = tag.get_text(strip=True)
                                    if txt and len(txt) > 10:
                                        text_blocks.append(txt)
                            if len(text_blocks) < 5:
                                for tag in soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6", "p", "li", "span", "div"]):
                                    txt = tag.get_text(strip=True)
                                    if txt and len(txt) > 20:
                                        text_blocks.append(txt)
                            content = f"[分頁: {title or link}]:\n" + "\n".join(text_blocks)
                            if content.strip():
                                all_contents.append(content)
                        except Exception as e:
                            self.logger.warning(f"[多分頁] 抓取 {link} 失敗: {e}")
                    await browser.close()
                    merged_content = "\n\n".join(all_contents)[:12000]  # 限制長度
                    if not merged_content or len(merged_content) < 100:
                        self.logger.warning("多分頁抓取後仍無有效內容")
                        return None
                    return merged_content
                else:
                    # 單頁網站維持原本策略
                    # 1. 嘗試等待常見內容 selector
                    selectors = [
                        'main', 'article', 'section', '.content', '.main', '.article', '#content', '#main', '#app', '#root'
                    ]
                    found = False
                    for sel in selectors:
                        try:
                            await page.wait_for_selector(sel, timeout=3000)
                            found = True
                            self.logger.info(f"等待 selector 成功: {sel}")
                            break
                        except Exception:
                            continue
                    # 2. 自動點擊展開/更多按鈕
                    expand_selectors = [
                        'button:has-text("展開")', 'button:has-text("更多")', 'button:has-text("Read more")',
                        'button:has-text("Show more")', 'a:has-text("展開")', 'a:has-text("更多")',
                        'a:has-text("Read more")', 'a:has-text("Show more")'
                    ]
                    for sel in expand_selectors:
                        try:
                            btns = await page.query_selector_all(sel)
                            for btn in btns:
                                await btn.click()
                                self.logger.info(f"自動點擊展開/更多按鈕: {sel}")
                                await page.wait_for_timeout(500)
                        except Exception:
                            continue
                    # 3. 滾動到底部，確保動態內容載入
                    await page.evaluate('window.scrollTo(0, document.body.scrollHeight)')
                    await page.wait_for_timeout(2000)
                    await page.evaluate('window.scrollTo(0, 0)')
                    await page.wait_for_timeout(1000)
                    # 4. 嘗試多種方式提取內容
                    html = await page.content()
                    soup = BeautifulSoup(html, "html.parser")
                    title = soup.title.string.strip() if soup.title else ""
                    meta_desc = soup.find("meta", attrs={"name": "description"})
                    desc = meta_desc["content"].strip() if meta_desc and meta_desc.get("content") else ""
                    text_blocks = []
                    main = soup.find("main") or soup.find("article") or soup.body
                    if main:
                        for tag in main.find_all(["h1", "h2", "h3", "h4", "h5", "h6", "p", "li"]):
                            txt = tag.get_text(strip=True)
                            if txt and len(txt) > 10:
                                text_blocks.append(txt)
                    if len(text_blocks) < 5:
                        for tag in soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6", "p", "li", "span", "div"]):
                            txt = tag.get_text(strip=True)
                            if txt and len(txt) > 20:
                                text_blocks.append(txt)
                    if len(text_blocks) < 3:
                        try:
                            inner_text = await page.evaluate('document.body.innerText')
                            if inner_text:
                                for line in inner_text.splitlines():
                                    line = line.strip()
                                    if len(line) > 20:
                                        text_blocks.append(line)
                                self.logger.info("已抓取 body.innerText")
                        except Exception as e:
                            self.logger.warning(f"抓取 body.innerText 失敗: {e}")
                    if len(text_blocks) < 3:
                        try:
                            text_content = await page.evaluate('document.body.textContent')
                            if text_content:
                                for line in text_content.splitlines():
                                    line = line.strip()
                                    if len(line) > 20:
                                        text_blocks.append(line)
                                self.logger.info("已抓取 body.textContent")
                        except Exception as e:
                            self.logger.warning(f"抓取 body.textContent 失敗: {e}")
                    await browser.close()
                    content = "\n".join(text_blocks)[:6000]  # 限制長度
                    if not content or len(content) < 100:
                        self.logger.warning("Playwright 沒有抓到有效內容，該網站可能需登入或有防爬蟲措施")
                        return None
                    return content
        except Exception as e:
            self.logger.error(f"Playwright內容擷取流程失敗: {str(e)}")
            return None


async def ocr_images_from_urls(image_urls: List[str]) -> str:
    """下載圖片並執行 OCR"""
    from pytesseract import pytesseract
    results = []

    for i, url in enumerate(image_urls):
        try:
            # 處理 base64 圖片
            if url.startswith('data:image/'):
                try:
                    import base64
                    # 取得 base64 編碼部分
                    header, encoded = url.split(",", 1)
                    img_data = base64.b64decode(encoded)
                    img = Image.open(BytesIO(img_data))
                except Exception as e:
                    logger.warning(f"❌ 無法解析 base64 圖片: {str(e)}")
                    continue
            # 處理一般 URL
            elif url.startswith("http"):
                response = requests.get(url)
                img = Image.open(BytesIO(response.content))
            else:
                logger.warning(f"❌ 不支援的圖片 URL 格式: {url[:100]}...")
                continue

            # 使用 UTF-8 編碼處理文字，添加 Tesseract 錯誤處理
            try:
                text = pytesseract.image_to_string(img, lang='eng')
                if text and text.strip():
                    # 確保文字使用 UTF-8 編碼
                    text_encoded = text.encode('utf-8', errors='ignore').decode('utf-8')
                    results.append(f"[Slide {i+1}]\n{text_encoded}")
            except Exception as ocr_e:
                # 處理 Tesseract 相關錯誤
                if "TesseractNotFoundError" in str(type(ocr_e)) or "tesseract" in str(ocr_e).lower():
                    logger.warning(f"⚠️ Tesseract OCR 不可用，跳過第 {i+1} 張圖片的文字提取: {str(ocr_e)}")
                    results.append(f"[Slide {i+1}]\n[OCR不可用 - 無法提取文字內容]")
                else:
                    # 其他 OCR 錯誤
                    logger.warning(f"❌ OCR 處理失敗 (第 {i+1} 張): {str(ocr_e)}")
                    results.append(f"[Slide {i+1}]\n[文字提取失敗]")
                
        except Exception as e:
            logger.warning(f"❌ 讀取圖片失敗: {str(e)}", exc_info=True)

    return "\n\n".join(results)

async def extract_company_name_from_message(message: str) -> Optional[str]:
    """從消息中提取公司名稱"""
    try:
        # 1. 檢查是否有明確的公司名稱標記
        company_patterns = [
            r'(?:^|\n|\s)(?:company|project|startup):\s*([^\n]+)',
            r'(?:^|\n|\s)(?:company|project|startup)\s*name:\s*([^\n]+)',
            r'(?:^|\n|\s)(?:company|project|startup)\s*name\s*is\s*([^\n]+)',
            r'(?:^|\n|\s)([A-Z][A-Za-z0-9\s&]+)\s*(?:is|are)\s*(?:a|an)\s*(?:company|startup|project)',
            r'(?:^|\n|\s)(?:about|introducing|presenting)\s+([A-Z][A-Za-z0-9\s&]+)',
            r'(?:^|\n|\s)([A-Z][A-Za-z0-9\s&]+)\s*(?:is|are)\s*(?:also|working|building)',
            r'(?:^|\n|\s)([A-Z][A-Za-z0-9\s&]+)\s*(?:is|are)\s*(?:governance|related)',
        ]
        
        for pattern in company_patterns:
            match = re.search(pattern, message, re.IGNORECASE | re.MULTILINE)
            if match:
                company_name = match.group(1).strip()
                # 清理公司名稱
                company_name = re.sub(r'\s+', ' ', company_name)  # 移除多餘空格
                company_name = re.sub(r'[^\w\s&-]', '', company_name)  # 只保留字母、數字、空格、&和-
                # 如果公司名稱太長，可能不是正確的匹配
                if len(company_name.split()) <= 3:  # 假設公司名稱不超過3個單詞
                    return company_name
        
        # 2. 檢查是否有 "about" 或 "introducing" 後面的名稱
        about_patterns = [
            r'(?:^|\n|\s)(?:about|introducing|presenting)\s+([A-Z][A-Za-z0-9\s&]+)',
            r'(?:^|\n|\s)([A-Z][A-Za-z0-9\s&]+)\s*(?:is|are)\s*(?:a|an)\s*(?:company|startup|project)',
            r'(?:^|\n|\s)([A-Z][A-Za-z0-9\s&]+)\s*(?:is|are)\s*(?:also|working|building)',
            r'(?:^|\n|\s)([A-Z][A-Za-z0-9\s&]+)\s*(?:is|are)\s*(?:governance|related)',
        ]
        
        for pattern in about_patterns:
            match = re.search(pattern, message, re.IGNORECASE | re.MULTILINE)
            if match:
                company_name = match.group(1).strip()
                company_name = re.sub(r'\s+', ' ', company_name)
                company_name = re.sub(r'[^\w\s&-]', '', company_name)
                if len(company_name.split()) <= 3:
                    return company_name
        
        # 3. 檢查是否有 "Blurb" 或 "Description" 後面的名稱
        blurb_patterns = [
            r'(?:^|\n|\s)(?:blurb|description):\s*([^\n]+)',
            r'(?:^|\n|\s)(?:blurb|description)\s*is\s*([^\n]+)',
        ]
        
        for pattern in blurb_patterns:
            match = re.search(pattern, message, re.IGNORECASE | re.MULTILINE)
            if match:
                text = match.group(1).strip()
                # 嘗試從描述中提取公司名稱
                company_match = re.search(r'([A-Z][A-Za-z0-9\s&]+)\s*(?:is|are)\s*(?:a|an)\s*(?:company|startup|project)', text)
                if company_match:
                    company_name = company_match.group(1).strip()
                    company_name = re.sub(r'\s+', ' ', company_name)
                    company_name = re.sub(r'[^\w\s&-]', '', company_name)
                    if len(company_name.split()) <= 3:
                        return company_name
        
        # 4. 特殊情況：直接查找大寫開頭的單詞
        words = message.split()
        for i, word in enumerate(words):
            if word[0].isupper() and len(word) > 1 and word.isalpha():
                # 檢查這個詞是否可能是公司名稱
                if i > 0 and words[i-1].lower() in ['about', 'introducing', 'presenting', 'company', 'project', 'startup']:
                    return word
                if i < len(words)-1 and words[i+1].lower() in ['is', 'are', 'working', 'building']:
                    return word
        
        return None
    except Exception as e:
        logger.error(f"提取公司名稱時出錯: {e}")
        return None

async def summarize_pitch_deck(ocr_text: str, message: str = "") -> Dict:
    """直接返回 OCR 文字結果"""
    # 確保輸入文字是 UTF-8 編碼
    if isinstance(ocr_text, str):
        ocr_text = ocr_text.encode('utf-8', errors='ignore').decode('utf-8')

    # 直接返回 OCR 結果
    return {
        "raw_content": ocr_text,
        "url": "",  # 添加空的 URL 字段以避免序列化問題
        "error": None  # 添加空的 error 字段以避免序列化問題
    }

async def debug_all_iframes(page):
    """
    儲存所有 iframe 的 HTML 並列出關鍵資訊（文字、圖片、canvas 數量）。
    """
    frames = page.frames
    print(f"\n🔍 檢測到 {len(frames)} 個 iframe。")

    for idx, frame in enumerate(frames):
        print(f"\n📦 解析 iframe [{idx}] URL: {frame.url}")

        try:
            # 儲存 HTML
            html = await frame.content()
            from utils.path_helper import PathHelper as _PH
            _ph = _PH()
            file_path = _ph.get(f"debug_frame_{idx}.html")
            with file_path.open("w", encoding="utf-8") as f:
                f.write(html)
            print(f"📝 已儲存 HTML 到 {file_path}")

            # 抓主要元素數量
            text_count = await frame.eval_on_selector_all('p', 'els => els.length')
            div_count = await frame.eval_on_selector_all('div', 'els => els.length')
            img_count = await frame.eval_on_selector_all('img', 'els => els.length')
            canvas_count = await frame.eval_on_selector_all('canvas', 'els => els.length')

            print(f"📊 內容統計：")
            print(f" - <p> 標籤：{text_count}")
            print(f" - <div> 標籤：{div_count}")
            print(f" - <img> 圖片：{img_count}")
            print(f" - <canvas> 元素：{canvas_count}")

            # 顯示可能性分析
            if canvas_count > 0:
                print("⚠️ 發現 canvas，可能是圖像渲染的 slide。")
            if text_count > 5 or div_count > 20:
                print("✅ 可能有文字內容。")
            if img_count > 5 and text_count < 2:
                print("🖼️ 很可能是純圖片型投影片。")

        except Exception as e:
            print(f"❌ 讀取 iframe {idx} 發生錯誤：{str(e)}")
    

#測試用驅動器
if __name__ == "__main__":

    async def main():
        message = """
        TrueNorth
        Crypto's first AI discovery engine that uses agentic technology to unlock a symbiotic user journey - from intent straight to outcome

        Co-founders

        Willy: Serial entrepreneur with a successful M&A exit (Series-B SaaS startup), Forbes 30 Under 30 China, ex-COO and acting CEO of WOO.

        Alex: PhD in AI & Domain-Specific Computing, ex-McKinsey, ex-Temasek, Head of Product, Strategy and Capital Market at Enflame (~USD3b pre-IPO AI chip startup), and the Tech Founding Partner of Iluvatar (~USD2b pre-IPO AI chip startup).

        Backed by
        Cyber Fund, Delphi Labs and founders, GPs from Layerzero, Virtuals, Selini, SEI, Merlin, Presto, LTP, Initial, Generative and more.

        Deck
        https://docsend.com/v/w7g7p/truenorth-pitch-deck-seed

        Website
        https://true-north.xyz/\
        """
        
        reader = DeckBrowser()
        try:
            await reader.initialize()
            results = await reader.process_input(message)
            if results:
                print(json.dumps(results, ensure_ascii=False, indent=2))
            else:
                print("⚠️ 沒有擷取到任何結果")
        except Exception as e:
            print(f"❌ 發生錯誤: {e}")
        finally:
            await reader.close()
            # 確保所有資源都被釋放
            import gc
            gc.collect()

    try:
        asyncio.run(main())
    except KeyboardInterrupt:
        print("Bot stopped by user.")
    except Exception as e:
        print(f"Fatal error: {e}")
    finally:
        # 確保所有資源都被釋放
        import gc
        gc.collect()